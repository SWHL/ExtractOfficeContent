# -*- encoding: utf-8 -*-
# @Author: SWHL
# @Contact: liekkaskono@163.com
from io import BytesIO
import argparse
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Union

import pandas as pd
import pptx
from pptx import Presentation

from .utils import mkdir


class ExtractPPT():
    def __init__(self, ):
        pass

    def __call__(self, ppt_content: Union[str, bytes],
                 save_img_dir: str = None) -> List:
        """Extract content and images of ppt.

        Args:
            ppt_content (str, bytes): the path of ppt.
            save_img_dir (str, optional): The directory for saving images. Defaults to None.

        Returns:
            List: txts from pptx.
        """
        if isinstance(ppt_content, str):
            if not Path(ppt_content).exists():
                raise FileNotFoundError(f'{ppt_content} does not exist.')
        elif isinstance(ppt_content, bytes):
            ppt_content = BytesIO(ppt_content)

        txts, imgs, charts = self.extract_all(ppt_content)
        if save_img_dir:
            if imgs:
                self.save_object(imgs, save_img_dir, suffix='png')

            if charts:
                self.save_object(charts, save_img_dir, suffix='xlsx')
        return list(txts.values())

    def extract_all(self, ppt_path: Union[str, bytes]) -> Tuple[Dict, Dict]:
        prs = Presentation(ppt_path)
        extract_txts, extract_imgs, extract_charts = {}, {}, {}
        for i, slide in enumerate(prs.slides):
            cur_page = i + 1
            cur_txts, cur_imgs, cur_charts = self.extract_one(slide)

            extract_txts[cur_page] = '\n'.join(cur_txts)
            extract_imgs.update({cur_page: cur_imgs})
            extract_charts.update({cur_page: cur_charts})
        return extract_txts, extract_imgs, extract_charts

    def extract_one(self, slide: pptx.slide.Slide) -> Tuple[List, List, List]:
        cur_page_content, cur_page_imgs, cur_page_charts = [], [], []
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = self.extract_text(shape.text)
                if txt:
                    cur_page_content.append(txt)
            elif shape.has_table:
                table_str = self.extract_table(shape.table)
                cur_page_content.append(table_str)
            elif shape.has_chart:
                excel_bytes = shape.chart.part.chart_workbook.xlsx_part.blob
                cur_page_charts.append(excel_bytes)
            elif hasattr(shape, 'image'):
                img_bytes = self.extract_image(shape.image)
                cur_page_imgs.append(img_bytes)
            else:
                pass
        return cur_page_content, cur_page_imgs, cur_page_charts

    @staticmethod
    def extract_text(shape_text: str) -> Optional[str]:
        txt = shape_text.strip()
        if txt:
            return txt
        return None

    @staticmethod
    def extract_table(table_value: pptx.table.Table) -> str:
        table_list = []
        for one_row in table_value.rows:
            each = ''
            for cell in one_row.cells:
                each += cell.text_frame.text + ','
            table_list.append(each)
        table_df = pd.DataFrame(table_list)
        return table_df.to_markdown(index=None)

    @staticmethod
    def extract_image(img_value: pptx.parts.image.Image) -> bytes:
        return img_value.blob

    @staticmethod
    def save_object(objs: Dict,
                    save_dir: Union[str, Path], suffix: str) -> None:
        mkdir(save_dir)
        for page_num, obj_list in objs.items():
            for i, img in enumerate(obj_list):
                save_full_path = Path(save_dir) / f'{page_num}_{i+1}.{suffix}'
                with open(str(save_full_path), 'wb') as f:
                    f.write(img)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('ppt_path', type=str)
    parser.add_argument('-img_dir', '--save_img_dir', type=str, default=None)
    args = parser.parse_args()

    ppt_extracter = ExtractPPT()
    res = ppt_extracter(args.ppt_path, save_img_dir=args.save_img_dir)
    print(res)


if __name__ == '__main__':
    main()
