# -*- encoding: utf-8 -*-
# @Author: SWHL
# @Contact: liekkaskono@163.com
from pathlib import Path
from typing import Dict, List, Tuple, Optional

import cv2
import numpy as np
import pandas as pd
import pptx
from pptx import Presentation

from .utils import mkdir, write_txt


class ExtractPPT():
    def __init__(self, ):
        pass

    def __call__(self, ppt_path: str,
                 is_save_to_txt: bool = False,
                 save_txt_dir: str = None,
                 is_save_img: bool = False,
                 save_img_dir: str = None) -> List:
        """Extract content and images of ppt.

        Args:
            ppt_path (str): the path of ppt.
            is_save_to_txt (bool, optional): Whether to save content to txt. Defaults to False.
            save_txt_dir (str, optional): The directory for saving txt. Defaults to None.
            is_save_img (bool, optional): Whether to save images to directory. Defaults to False.
            save_img_dir (str, optional): The directory for saving images. Defaults to None.

        Returns:
            List: txts from pptx.
        """
        if is_save_to_txt and save_txt_dir is None:
            raise ValueError(
                'When is_save_to_txt is True, save_txt_dir must not be None.')

        if is_save_img and save_img_dir is None:
            raise ValueError(
                'When is_save_img is True, save_img_dir must be not None.')

        txts, imgs = self.extract_all(ppt_path)

        if is_save_to_txt and save_txt_dir:
            mkdir(save_txt_dir)
            full_txt_path = Path(save_txt_dir) / f'{Path(ppt_path).stem}.txt'
            write_txt(full_txt_path, list(txts.values()))

        if is_save_img and save_img_dir:
            mkdir(save_img_dir)
            for page_num, img_list in imgs.items():
                for i, img in enumerate(img_list):
                    save_full_path = Path(save_img_dir) / \
                        f'{page_num}_{i+1}.png'
                    cv2.imwrite(str(save_full_path), img)
        return list(txts.values())

    def extract_all(self, ppt_path: str) -> Tuple[Dict, Dict]:
        prs = Presentation(ppt_path)
        extract_txts, extract_imgs = {}, {}
        for i, slide in enumerate(prs.slides):
            cur_page = i + 1
            cur_txts, cur_imgs = self.extract_one(slide)

            extract_txts[cur_page] = '\n'.join(cur_txts)
            for cur_img in cur_imgs:
                extract_imgs.setdefault(cur_page, []).append(cur_img)
        return extract_txts, extract_imgs

    def extract_one(self, slide: pptx.slide.Slide) -> List:
        cur_page_content, cur_page_imgs = [], []
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = self.extract_text(shape.text)
                if txt:
                    cur_page_content.append(txt)
            elif shape.has_table:
                table_str = self.extract_table(shape.table)
                cur_page_content.append(table_str)
            elif shape.has_chart:
                pass
            elif hasattr(shape, 'image'):
                img = self.extract_image(shape.image)
                cur_page_imgs.append(img)
            else:
                pass
        return cur_page_content, cur_page_imgs

    @staticmethod
    def extract_text(shape_text: str) -> Optional[str]:
        txt = shape_text.strip()
        if txt:
            return txt
        return None

    @staticmethod
    def extract_table(table_value: pptx.table.Table) -> str:
        table_list = []
        for one_row in table_value.rows:
            each = ''
            for cell in one_row.cells:
                each += cell.text_frame.text + ','
            table_list.append(each)
        table_df = pd.DataFrame(table_list)
        return table_df.to_string()

    @staticmethod
    def extract_image(img_value: pptx.parts.image.Image) -> np.ndarray:
        img_blob = img_value.blob
        img_np = np.frombuffer(img_blob, dtype=np.uint8)
        img_array = cv2.imdecode(img_np, cv2.IMREAD_COLOR)
        return img_array
