# -*- encoding: utf-8 -*-
# @Author: SWHL
# @Contact: liekkaskono@163.com
from pathlib import Path

import pandas as pd
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pptx
from pptx import Presentation
from utils import mkdir, write_txt


class ExtractPPTText():
    def __init__(self, ppt_path: str, save_dir: str):
        self.prs = Presentation(ppt_path)
        self.save_dir = Path(save_dir) / Path(ppt_path).stem
        mkdir(self.save_dir)

    def __call__(self,):
        for i, slide in enumerate(self.prs.slides):
            print(i)
            cur_page_content = []
            for shape in slide.shapes:
                # if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                #     pass

                if shape.has_text_frame:
                    txt = self.extract_text(shape.text)
                    if txt:
                        cur_page_content.append(txt)
                elif shape.has_table:
                    table_str = self.extract_table(shape.table)
                    cur_page_content.append(table_str)
                elif shape.has_chart:
                    pass
                elif hasattr(shape, 'image'):
                    self.save_image(shape.image, page_num=i)
                else:
                    pass

            cur_page_path = self.save_dir / f'{i:0>2}.txt'
            write_txt(cur_page_path, cur_page_content)

    @staticmethod
    def extract_text(shape_text):
        txt = shape_text.strip()
        if txt:
            return txt
        return None

    @staticmethod
    def extract_table(table_value: pptx.table.Table) -> str:
        table_list = []
        for one_row in table_value.rows:
            each = ''
            for cell in one_row.cells:
                each += cell.text_frame.text + ','
            table_list.append(each)
        table_df = pd.DataFrame(table_list)
        return table_df.to_string()

    def save_image(self, img_value, page_num: int):
        img_name = img_value.filename

        save_img_dir = self.save_dir / 'images'
        mkdir(save_img_dir)
        save_img_path = save_img_dir / f'{page_num}_{img_name}'

        img_blob = img_value.blob
        with open(save_img_path, "wb") as f:
            f.write(img_blob)


if __name__ == '__main__':
    ppt_path = 'test_files/简约活动策划方案汇报PPT模板.pptx'

    ppt_extracter = ExtractPPTText(ppt_path, save_dir='outputs')

    res = ppt_extracter()
    print(res)
